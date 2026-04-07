"""
GenTwin Hackathon Presentation Generator
Creates a modern 8-slide PowerPoint deck with blue primary color
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# Modern blue primary color
PRIMARY_COLOR = RGBColor(37, 99, 235)  # #2563EB (blue-600)
DARK_GRAY = RGBColor(31, 41, 55)  # #1F2937
LIGHT_GRAY = RGBColor(107, 114, 128)  # #6B7280
WHITE = RGBColor(255, 255, 255)

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)  # 16:9 aspect ratio
    prs.slide_height = Inches(5.625)
    
    # Slide 1: Solution Overview (Title + Intro + Objectives)
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Title
    title = slide1.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(9), Inches(1))
    tf = title.text_frame
    tf.text = "GenTwin"
    p = tf.paragraphs[0]
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle = slide1.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9), Inches(0.5))
    tf = subtitle.text_frame
    tf.text = "AI-Powered Cybersecurity Digital Twin for Industrial Water Treatment"
    p = tf.paragraphs[0]
    p.font.size = Pt(20)
    p.font.color.rgb = DARK_GRAY
    p.alignment = PP_ALIGN.CENTER
    
    # Introduction
    intro_box = slide1.shapes.add_textbox(Inches(0.8), Inches(2.8), Inches(4), Inches(2))
    tf = intro_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Introduction"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    p.space_after = Pt(8)
    
    intro_text = """Real-time cyber threat detection and mitigation for critical infrastructure using advanced AI, digital twin simulation, and explainable ML."""
    p = tf.add_paragraph()
    p.text = intro_text
    p.font.size = Pt(14)
    p.font.color.rgb = DARK_GRAY
    p.line_spacing = 1.3
    
    # Objectives
    obj_box = slide1.shapes.add_textbox(Inches(5.2), Inches(2.8), Inches(4), Inches(2))
    tf = obj_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Objectives"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    p.space_after = Pt(8)
    
    objectives = [
        "Detect novel cyber attacks in real-time",
        "Simulate attack impact on plant safety",
        "Provide explainable threat insights",
        "Enable adaptive defense strategies"
    ]
    
    for obj in objectives:
        p = tf.add_paragraph()
        p.text = f"• {obj}"
        p.font.size = Pt(13)
        p.font.color.rgb = DARK_GRAY
        p.level = 0
        p.space_before = Pt(4)
    
    # Slide 2: System Architecture
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = title.text_frame
    tf.text = "System Architecture"
    p = tf.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    
    # Architecture layers
    arch_data = [
        ("Layer 1: Data Collection", "SWaT sensor data (51 channels)\nReal-time virtual sensors\nNormal & attack datasets"),
        ("Layer 2: ML Models", "VAE: Anomaly detection\nLSTM-AE: Temporal patterns\nCGAN: Attack generation"),
        ("Layer 3: Digital Twin", "SimPy plant simulation\nPhysical impact analysis\nSafety constraint validation"),
        ("Layer 4: Analysis", "Security gap discovery\nSHAP/LIME explainability\nBlindspot scoring")
    ]
    
    y_pos = 1.2
    for layer, desc in arch_data:
        # Layer box
        box = slide2.shapes.add_shape(
            1,  # Rectangle
            Inches(0.8), Inches(y_pos),
            Inches(8.4), Inches(0.85)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = PRIMARY_COLOR
        box.line.width = Pt(2)
        
        # Layer title
        tf = box.text_frame
        tf.margin_top = Inches(0.05)
        tf.margin_left = Inches(0.15)
        p = tf.paragraphs[0]
        p.text = layer
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_COLOR
        
        # Description
        p = tf.add_paragraph()
        p.text = desc
        p.font.size = Pt(11)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(4)
        
        y_pos += 1.0
    
    # Slide 3: Tech Stack
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = title.text_frame
    tf.text = "Technology Stack"
    p = tf.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    
    # Tech categories
    tech_data = [
        ("Deep Learning", ["PyTorch 2.0+", "VAE & LSTM Autoencoders", "Conditional GAN"]),
        ("Explainability", ["SHAP (KernelExplainer)", "LIME (Tabular)", "Feature Attribution"]),
        ("Simulation", ["SimPy 4.0+", "Discrete Event Simulation", "Safety Violation Modeling"]),
        ("Data Science", ["Pandas & NumPy", "Scikit-learn", "NetworkX (Graph Analysis)"]),
        ("Visualization", ["Matplotlib & Seaborn", "Plotly (Interactive)", "Real-time Dashboards"]),
        ("Innovation", ["RL Q-Learning", "Immunity Scoring", "DNA Fingerprinting"])
    ]
    
    x_positions = [0.5, 3.5, 6.5]
    y_positions = [1.2, 3.2]
    
    for idx, (category, items) in enumerate(tech_data):
        row = idx // 3
        col = idx % 3
        
        x = x_positions[col]
        y = y_positions[row]
        
        # Category box
        box = slide3.shapes.add_shape(
            1,  # Rectangle
            Inches(x), Inches(y),
            Inches(2.8), Inches(1.6)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = PRIMARY_COLOR
        box.line.width = Pt(2)
        
        tf = box.text_frame
        tf.margin_top = Inches(0.1)
        tf.margin_left = Inches(0.15)
        tf.word_wrap = True
        
        # Category title
        p = tf.paragraphs[0]
        p.text = category
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_COLOR
        p.space_after = Pt(6)
        
        # Items
        for item in items:
            p = tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(10)
            p.font.color.rgb = DARK_GRAY
            p.space_before = Pt(2)
    
    # Slide 4: Innovations
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title = slide4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = title.text_frame
    tf.text = "Core Innovations"
    p = tf.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    
    innovations = [
        {
            "title": "🎯 RL Adaptive Defense",
            "desc": "Q-learning agent that learns optimal mitigation strategies (monitor, alert, isolate, shutdown) based on real-time threat severity levels."
        },
        {
            "title": "🛡️ Immunity System",
            "desc": "Per-stage resilience scoring combining detection capability, impact resistance, and gap density to identify vulnerable plant sections."
        },
        {
            "title": "🧬 Cyber DNA Fingerprinting",
            "desc": "SHA-256 attack signatures with statistical features for rapid threat identification, deduplication, and similarity-based retrieval."
        },
        {
            "title": "⏱️ Predictive Timeline",
            "desc": "Chronological incident reconstruction with impact scores, violation counts, and critical gap markers for forensic analysis."
        }
    ]
    
    y_pos = 1.2
    for innov in innovations:
        # Innovation box
        box = slide4.shapes.add_shape(
            1,  # Rectangle
            Inches(0.8), Inches(y_pos),
            Inches(8.4), Inches(0.9)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = PRIMARY_COLOR
        box.line.width = Pt(2)
        
        tf = box.text_frame
        tf.margin_top = Inches(0.08)
        tf.margin_left = Inches(0.15)
        tf.word_wrap = True
        
        # Title
        p = tf.paragraphs[0]
        p.text = innov["title"]
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_COLOR
        
        # Description
        p = tf.add_paragraph()
        p.text = innov["desc"]
        p.font.size = Pt(11)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(4)
        
        y_pos += 1.0
    
    # Slide 5: Flowchart
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title = slide5.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = title.text_frame
    tf.text = "Pipeline Flowchart"
    p = tf.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    
    # Flow steps
    flow_steps = [
        "Data Ingestion\nSWaT Sensors (51 channels)",
        "Preprocessing\nNormalization & Sequencing",
        "ML Detection\nVAE + LSTM-AE Anomaly Scores",
        "Attack Generation\nCGAN Synthetic Threats",
        "SimPy Simulation\nPhysical Impact Analysis",
        "Gap Discovery\nHigh Impact + Low Detection",
        "Explainability\nSHAP + LIME Insights",
        "Adaptive Defense\nRL Policy + Timeline"
    ]
    
    x_start = 0.8
    box_width = 1.9
    gap = 0.3
    y_pos = 1.3
    
    for idx, step in enumerate(flow_steps):
        row = idx // 4
        col = idx % 4
        
        x = x_start + col * (box_width + gap)
        y = y_pos + row * 1.5
        
        # Flow box
        box = slide5.shapes.add_shape(
            1,  # Rectangle
            Inches(x), Inches(y),
            Inches(box_width), Inches(1.1)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = PRIMARY_COLOR if idx % 2 == 0 else WHITE
        box.line.color.rgb = PRIMARY_COLOR
        box.line.width = Pt(2)
        
        tf = box.text_frame
        tf.text = step
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = WHITE if idx % 2 == 0 else PRIMARY_COLOR
        
        # Arrow (except last)
        if col < 3 and idx < len(flow_steps) - 1:
            arrow = slide5.shapes.add_shape(
                5,  # Arrow
                Inches(x + box_width), Inches(y + 0.5),
                Inches(gap), Inches(0.1)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = PRIMARY_COLOR
            arrow.line.color.rgb = PRIMARY_COLOR
    
    # Slide 6: Real-world Impact
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title = slide6.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = title.text_frame
    tf.text = "Real-World Impact"
    p = tf.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    
    impact_areas = [
        {
            "icon": "🏭",
            "title": "Critical Infrastructure Protection",
            "points": [
                "Water treatment plants",
                "Chemical processing facilities",
                "Power generation systems",
                "Oil & gas operations"
            ]
        },
        {
            "icon": "🔒",
            "title": "Cyber-Physical Security",
            "points": [
                "Real-time threat detection",
                "Safety violation prevention",
                "Attack impact mitigation",
                "Compliance monitoring"
            ]
        },
        {
            "icon": "💡",
            "title": "Operational Benefits",
            "points": [
                "Reduced downtime",
                "Lower insurance costs",
                "Regulatory compliance",
                "Improved incident response"
            ]
        }
    ]
    
    x_positions = [0.5, 3.5, 6.5]
    for idx, area in enumerate(impact_areas):
        x = x_positions[idx]
        
        # Area box
        box = slide6.shapes.add_shape(
            1,  # Rectangle
            Inches(x), Inches(1.2),
            Inches(2.8), Inches(3.8)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = PRIMARY_COLOR
        box.line.width = Pt(2)
        
        tf = box.text_frame
        tf.margin_top = Inches(0.15)
        tf.margin_left = Inches(0.15)
        tf.word_wrap = True
        
        # Icon + Title
        p = tf.paragraphs[0]
        p.text = f"{area['icon']} {area['title']}"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_COLOR
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(10)
        
        # Points
        for point in area["points"]:
            p = tf.add_paragraph()
            p.text = f"• {point}"
            p.font.size = Pt(11)
            p.font.color.rgb = DARK_GRAY
            p.space_before = Pt(6)
    
    # Slide 7: Results
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title = slide7.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = title.text_frame
    tf.text = "Results & Performance"
    p = tf.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    
    # Metrics
    metrics = [
        ("1,500+", "Synthetic Attack Vectors Generated", "Via CGAN for comprehensive testing"),
        ("300", "Simulated Attack Scenarios", "With physical impact assessment"),
        ("50+", "Critical Security Gaps Identified", "High impact, low detection attacks"),
        ("7", "Analytical Artifacts", "Complete threat intelligence pipeline"),
        ("6", "Plant Stages Analyzed", "Full SWaT process coverage"),
        ("51", "Sensor Channels Monitored", "Real-time anomaly detection")
    ]
    
    x_positions = [0.8, 3.5, 6.2]
    y_positions = [1.3, 3.5]
    
    for idx, (value, label, detail) in enumerate(metrics):
        row = idx // 3
        col = idx % 3
        
        x = x_positions[col]
        y = y_positions[row]
        
        # Metric box
        box = slide7.shapes.add_shape(
            1,  # Rectangle
            Inches(x), Inches(y),
            Inches(2.5), Inches(1.7)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = PRIMARY_COLOR if idx < 3 else WHITE
        box.line.color.rgb = PRIMARY_COLOR
        box.line.width = Pt(2)
        
        tf = box.text_frame
        tf.margin_top = Inches(0.15)
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        # Value
        p = tf.paragraphs[0]
        p.text = value
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = WHITE if idx < 3 else PRIMARY_COLOR
        p.alignment = PP_ALIGN.CENTER
        
        # Label
        p = tf.add_paragraph()
        p.text = label
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = WHITE if idx < 3 else PRIMARY_COLOR
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(6)
        
        # Detail
        p = tf.add_paragraph()
        p.text = detail
        p.font.size = Pt(9)
        p.font.color.rgb = WHITE if idx < 3 else LIGHT_GRAY
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(4)
    
    # Slide 8: Feature Scope
    slide8 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title = slide8.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = title.text_frame
    tf.text = "Feature Scope"
    p = tf.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    
    # Current Features
    current_box = slide8.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(4), Inches(2))
    tf = current_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "✅ Current Features (Person 1 + 2)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    p.space_after = Pt(8)
    
    current_features = [
        "VAE & LSTM-AE anomaly detection",
        "CGAN attack generation (6 stages)",
        "SimPy digital twin simulation",
        "SHAP/LIME explainability",
        "Security gap discovery",
        "RL adaptive defense (Q-learning)",
        "Immunity scoring system",
        "Cyber DNA fingerprinting",
        "Incident timeline builder"
    ]
    
    for feat in current_features:
        p = tf.add_paragraph()
        p.text = f"• {feat}"
        p.font.size = Pt(11)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(3)
    
    # Future Roadmap
    future_box = slide8.shapes.add_textbox(Inches(5.2), Inches(1.2), Inches(4), Inches(2))
    tf = future_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "🚀 Future Roadmap (Person 3+)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    p.space_after = Pt(8)
    
    future_features = [
        "Real-time prediction pipeline",
        "Advanced autoencoder refinement",
        "Multi-model ensemble",
        "Live sensor integration",
        "Interactive dashboard",
        "Alert notification system",
        "Automated response triggers",
        "Cloud deployment",
        "Multi-site scaling"
    ]
    
    for feat in future_features:
        p = tf.add_paragraph()
        p.text = f"• {feat}"
        p.font.size = Pt(11)
        p.font.color.rgb = LIGHT_GRAY
        p.space_before = Pt(3)
    
    # Key Metrics Footer
    footer = slide8.shapes.add_textbox(Inches(0.8), Inches(3.8), Inches(8.4), Inches(1.2))
    tf = footer.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "📊 Project Metrics"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    p.space_after = Pt(6)
    
    metrics_text = "✓ 74/74 pipeline tests passed  |  ✓ 10+ Python modules  |  ✓ 7 analytical artifacts  |  ✓ Production-ready codebase  |  ✓ Comprehensive documentation"
    p = tf.add_paragraph()
    p.text = metrics_text
    p.font.size = Pt(10)
    p.font.color.rgb = DARK_GRAY
    
    # Save presentation
    output_path = "GenTwin_Hackathon_Presentation.pptx"
    prs.save(output_path)
    print(f"✅ Presentation created: {output_path}")
    return output_path

if __name__ == "__main__":
    create_presentation()
